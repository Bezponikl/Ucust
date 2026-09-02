"""
test_document_collector_and_rag.py
======================================================================
Тестирование:
1. DocumentCollector: Извлечение текста и таблиц из PDF, DOCX, PPTX.
2. Онбординг с прикреплением файлов клиента (PDF, DOCX, PPTX).
3. Индексация документов в Clean RAG Pipeline (категория client_files).
4. Генерация постов с извлечением фактов и прайса из загруженных файлов.
======================================================================
"""

import sys
if hasattr(sys.stdout, 'reconfigure'):
    sys.stdout.reconfigure(encoding='utf-8')

import os
import asyncio
from datetime import datetime
from sqlalchemy import create_engine
from sqlalchemy.orm import sessionmaker

AI_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
if AI_ROOT not in sys.path:
    sys.path.insert(0, AI_ROOT)

from storage.db import Base
from storage.models import UserProfile
from collectors.document_collector import DocumentCollector
from core.orchestrator import UnifiedOrchestrator


def create_sample_pdf(file_path: str):
    """Создает тестовый PDF-прайс."""
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas
    c = canvas.Canvas(file_path, pagesize=letter)
    c.setFont("Helvetica", 14)
    c.drawString(50, 750, "AutoPro Detailing - Official Price List 2026")
    c.setFont("Helvetica", 11)
    c.drawString(50, 720, "1. Ceramic Coating 9H Pro: 45,000 RUB (3 years warranty)")
    c.drawString(50, 700, "2. Full Body PPF Film: 180,000 RUB (Self-healing premium film)")
    c.drawString(50, 680, "3. Interior Deep Cleaning: 15,000 RUB (Ozone sterilization included)")
    c.drawString(50, 650, "Special USP: Free car wash every 3 months for all ceramic package owners.")
    c.showPage()
    c.save()


def create_sample_docx(file_path: str):
    """Создает тестовый DOCX с описанием услуг и гарантией."""
    import docx
    doc = docx.Document()
    doc.add_heading("Коммерческое предложение AutoPro Detailing", level=1)
    doc.add_paragraph("Мы используем только немецкие премиальные составы Koch Chemie и Gyeon.")
    doc.add_paragraph("Все мастера прошли сертификацию в Мюнхене и имеют опыт от 7 лет.")
    
    # Таблица гарантий
    table = doc.add_table(rows=1, cols=3)
    hdr_cells = table.rows[0].cells
    hdr_cells[0].text = "Услуга"
    hdr_cells[1].text = "Гарантия"
    hdr_cells[2].text = "Срок выполнения"
    
    row_cells = table.add_row().cells
    row_cells[0].text = "Керамика 9H"
    row_cells[1].text = "36 месяцев"
    row_cells[2].text = "2 дня"
    
    doc.save(file_path)


def create_sample_pptx(file_path: str):
    """Создает тестовую презентацию PPTX."""
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "AutoPro Detailing Презентация"
    slide.placeholders[1].text = "Премиальный уход за автомобилями в Санкт-Петербурге\nАдрес: ул. Профессора Попова, 23\nТелефон: +7 (812) 999-88-77"
    prs.save(file_path)


async def run_document_collector_test():
    print("=" * 80)
    print("📄 ТЕСТ МОДУЛЯ DOCUMENT COLLECTOR (PDF, DOCX, PPTX) И ИНДЕКСАЦИИ В RAG")
    print("=" * 80)

    scratch_dir = os.path.join(AI_ROOT, "temp_test_docs")
    os.makedirs(scratch_dir, exist_ok=True)

    pdf_file = os.path.join(scratch_dir, "price_list.pdf")
    docx_file = os.path.join(scratch_dir, "commercial_offer.docx")
    pptx_file = os.path.join(scratch_dir, "pitch_deck.pptx")

    create_sample_pdf(pdf_file)
    create_sample_docx(docx_file)
    create_sample_pptx(pptx_file)

    # ------------------------------------------------------------------
    # ЭТАП 1: ПРОВЕРКА ИЗВЛЕЧЕНИЯ ИЗ PDF, DOCX, PPTX
    # ------------------------------------------------------------------
    print("\n--- [ЭТАП 1] Прямое извлечение через DocumentCollector ---")
    collector = DocumentCollector()

    docs_res = collector.extract_documents_batch([pdf_file, docx_file, pptx_file])
    assert len(docs_res) == 3

    for d in docs_res:
        print(f"✅ Файл: {d['file_name']} [{d['format'].upper()}] — {d['char_count']} символов, {d['units_count']} {d['unit_type']}")
        assert d["status"] == "success"
        assert d["char_count"] > 20

    assert "Ceramic Coating 9H Pro" in docs_res[0]["raw_text"]
    assert "Koch Chemie" in docs_res[1]["raw_text"]
    assert "Профессора Попова" in docs_res[2]["raw_text"]

    # ------------------------------------------------------------------
    # ЭТАП 2: ОНБОРДИНГ С ФАЙЛАМИ И ИНДЕКСАЦИЕЙ В RAG
    # ------------------------------------------------------------------
    print("\n--- [ЭТАП 2] Онбординг компании с прикреплением файлов ---")
    db_path = os.path.join(AI_ROOT, "test_docs_db.db")
    if os.path.exists(db_path):
        try:
            os.remove(db_path)
        except Exception:
            pass

    engine_db = create_engine(f"sqlite:///{db_path}")
    Base.metadata.create_all(engine_db)
    SessionLocal = sessionmaker(bind=engine_db)
    db_session = SessionLocal()

    orchestrator = UnifiedOrchestrator(db_session=db_session)

    onboard_res = await orchestrator.execute_task(
        task_type="onboard_user",
        user_data={
            "user_id": "user_doc_client_44",
            "company_name": "AutoPro Detailing",
            "niche": "Премиальный детейлинг",
            "city": "Санкт-Петербург",
            "country": "Россия",
            "documents": [pdf_file, docx_file, pptx_file],
            "answers": {
                "step1": {"brand_name": "AutoPro Detailing", "tone": "Премиальный и экспертный"},
                "step2": {"target_audience": "Владельцы авто премиум-класса"},
                "step3": {"advantages": "Немецкая химия, сертифицированные мастера"},
                "step4": {"services": [{"title": "Керамика 9H Pro", "price": "45 000 руб."}]},
                "step5": {"goals": ["Увеличение чека"]}
            }
        },
        session_id="sess_onboard_docs"
    )

    assert onboard_res.get("status") == "success"
    profile_id = onboard_res.get("profile_id")
    print(f"✅ Онбординг завершен! Профиль ID: {profile_id}")

    saved_profile = db_session.query(UserProfile).filter(UserProfile.id == profile_id).first()
    assert saved_profile is not None
    assert saved_profile.brand_dossier.get("documents_dossier") is not None
    print(f"✅ В SQL brand_dossier сохранен documents_dossier длиной {len(saved_profile.brand_dossier['documents_dossier'])} симв.")

    # ------------------------------------------------------------------
    # ЭТАП 3: ГЕНЕРАЦИЯ ПОСТА С ИЗВЛЕЧЕНИЕМ ФАКТОВ ИЗ ФАЙЛОВ ИЗ RAG
    # ------------------------------------------------------------------
    print("\n--- [ЭТАП 3] Генерация поста с обращением к знаниям из файлов ---")
    post_res = await orchestrator.execute_task(
        task_type="generate_post",
        user_data={
            "profile_id": profile_id,
            "prompt": "Пост о гарантиях на керамику 9H Pro и немецкой химии Koch Chemie из нашего коммерческого предложения",
            "generate_image": True
        },
        session_id="sess_gen_doc_post"
    )

    assert post_res.get("status") == "success"
    print("\n✍️ Сгенерированный пост на базе файлов клиента:")
    print("-" * 50)
    print(post_res["post_text"])
    print("-" * 50)
    print(f"📸 Фото: {post_res.get('photo_path')}")

    # Очистка
    db_session.close()
    engine_db.dispose()
    try:
        if os.path.exists(db_path):
            os.remove(db_path)
        for f in [pdf_file, docx_file, pptx_file]:
            if os.path.exists(f):
                os.remove(f)
        if os.path.exists(scratch_dir):
            os.rmdir(scratch_dir)
    except Exception:
        pass

    print("\n" + "=" * 80)
    print("🎉 ВСЕ ТЕСТЫ DOCUMENT COLLECTOR И RAG-ИНДЕКСАЦИИ УСПЕШНО ПРОЙДЕНЫ!")
    print("=" * 80)


if __name__ == "__main__":
    asyncio.run(run_document_collector_test())
