"""
DocumentCollector for UCust.AI.
Извлечение текста, таблиц и структуры из файлов клиентов (PDF, DOCX, PPTX, TXT)
с автоматической очисткой, семантическим анализом и индексацией в Clean RAG.
"""

from __future__ import annotations

import os
import re
import logging
from typing import List, Dict, Any, Optional

logger = logging.getLogger("document_collector")


class DocumentCollector:
    """
    Коллектор документов клиентов:
    - Извлекает контент из PDF (презентации, КП), DOCX (договоры, прайсы) и PPTX (маркетинг-киты).
    - Очищает текст от артефактов верстки и колонтитулов.
    - Семантически обогащает RAG-память бренда документами категории 'client_files'.
    """

    SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".txt", ".md", ".csv"}

    def __init__(self):
        pass

    def extract_text_from_file(self, file_path: str) -> Dict[str, Any]:
        """
        Извлекает сырой и очищенный текст из одного файла любого поддерживаемого формата.
        """
        if not os.path.exists(file_path):
            return {
                "status": "error",
                "file_path": file_path,
                "error": f"Файл не найден: {file_path}"
            }

        file_name = os.path.basename(file_path)
        ext = os.path.splitext(file_name)[1].lower()

        if ext not in self.SUPPORTED_EXTENSIONS:
            return {
                "status": "error",
                "file_path": file_path,
                "error": f"Формат '{ext}' не поддерживается. Разрешены: {', '.join(self.SUPPORTED_EXTENSIONS)}"
            }

        try:
            if ext == ".pdf":
                return self._extract_pdf(file_path, file_name)
            elif ext == ".docx":
                return self._extract_docx(file_path, file_name)
            elif ext == ".pptx":
                return self._extract_pptx(file_path, file_name)
            else:
                return self._extract_plain_text(file_path, file_name)
        except Exception as e:
            logger.error(f"[DocumentCollector] ❌ Ошибка парсинга {file_name}: {e}")
            return {
                "status": "error",
                "file_path": file_path,
                "file_name": file_name,
                "error": str(e)
            }

    def _clean_text(self, text: str) -> str:
        """Очищает текст от мусора, повторяющихся пробелов и разрывов строк."""
        if not text:
            return ""
        # Удаляем специфичные непечатные символы
        text = text.replace('\x00', '').replace('\ufeff', '')
        # Убираем множественные пустые строки
        text = re.sub(r'\n{3,}', '\n\n', text)
        # Убираем множественные пробелы
        text = re.sub(r'[ \t]{2,}', ' ', text)
        return text.strip()

    def _extract_pdf(self, file_path: str, file_name: str) -> Dict[str, Any]:
        """Извлечение текста из PDF через pypdf."""
        import pypdf
        reader = pypdf.PdfReader(file_path)
        num_pages = len(reader.pages)
        pages_content = []

        for idx, page in enumerate(reader.pages):
            page_text = page.extract_text() or ""
            cleaned = self._clean_text(page_text)
            if cleaned:
                pages_content.append(f"--- Страница {idx + 1} ---\n{cleaned}")

        full_text = "\n\n".join(pages_content)
        return {
            "status": "success",
            "format": "pdf",
            "file_name": file_name,
            "file_path": file_path,
            "units_count": num_pages,
            "unit_type": "pages",
            "raw_text": full_text,
            "char_count": len(full_text)
        }

    def _extract_docx(self, file_path: str, file_name: str) -> Dict[str, Any]:
        """Извлечение текста и таблиц из DOCX через python-docx."""
        import docx
        doc = docx.Document(file_path)
        paragraphs_text = []

        # 1. Параграфы и заголовки
        for p in doc.paragraphs:
            t = p.text.strip()
            if t:
                # Если заголовок, выделяем
                if p.style and "heading" in p.style.name.lower():
                    paragraphs_text.append(f"\n### {t}\n")
                else:
                    paragraphs_text.append(t)

        # 2. Таблицы (прайсы, списки услуг)
        for t_idx, table in enumerate(doc.tables):
            table_rows = []
            for row in table.rows:
                row_cells = [c.text.strip() for c in row.cells if c.text.strip()]
                if row_cells:
                    table_rows.append(" | ".join(row_cells))
            if table_rows:
                paragraphs_text.append(f"\n[Таблица {t_idx + 1}]:\n" + "\n".join(table_rows))

        full_text = self._clean_text("\n".join(paragraphs_text))
        return {
            "status": "success",
            "format": "docx",
            "file_name": file_name,
            "file_path": file_path,
            "units_count": len(doc.paragraphs),
            "unit_type": "paragraphs",
            "raw_text": full_text,
            "char_count": len(full_text)
        }

    def _extract_pptx(self, file_path: str, file_name: str) -> Dict[str, Any]:
        """Извлечение слайдов, надписей и таблиц из презентаций PPTX."""
        from pptx import Presentation
        prs = Presentation(file_path)
        slides_content = []

        for idx, slide in enumerate(prs.slides):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        txt = paragraph.text.strip()
                        if txt:
                            slide_texts.append(txt)
                elif shape.has_table:
                    for row in shape.table.rows:
                        row_cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if row_cells:
                            slide_texts.append(" | ".join(row_cells))

            if slide_texts:
                cleaned_slide = self._clean_text("\n".join(slide_texts))
                slides_content.append(f"--- Слайд {idx + 1} ---\n{cleaned_slide}")

        full_text = "\n\n".join(slides_content)
        return {
            "status": "success",
            "format": "pptx",
            "file_name": file_name,
            "file_path": file_path,
            "units_count": len(prs.slides),
            "unit_type": "slides",
            "raw_text": full_text,
            "char_count": len(full_text)
        }

    def _extract_plain_text(self, file_path: str, file_name: str) -> Dict[str, Any]:
        """Извлечение текста из TXT, MD, CSV."""
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            content = f.read()
        cleaned = self._clean_text(content)
        return {
            "status": "success",
            "format": os.path.splitext(file_name)[1].replace(".", ""),
            "file_name": file_name,
            "file_path": file_path,
            "units_count": 1,
            "unit_type": "file",
            "raw_text": cleaned,
            "char_count": len(cleaned)
        }

    def extract_documents_batch(self, file_paths: List[str]) -> List[Dict[str, Any]]:
        """Пакетное извлечение текстов из нескольких файлов."""
        results = []
        for fp in file_paths:
            res = self.extract_text_from_file(fp)
            results.append(res)
        return results

    async def sync_documents_to_rag(
        self,
        company_name: str,
        niche: str,
        extracted_docs: List[Dict[str, Any]],
        rag_pipeline: Any
    ) -> int:
        """
        Индексирует извлеченные материалы из файлов в Clean RAG Pipeline
        в семантическую категорию 'client_files'.
        """
        if not rag_pipeline or not extracted_docs:
            return 0

        from rag.models import Document
        docs_to_ingest = []

        for idx, doc_item in enumerate(extracted_docs):
            if doc_item.get("status") != "success" or not doc_item.get("raw_text"):
                continue

            file_name = doc_item.get("file_name", f"doc_{idx+1}")
            fmt = doc_item.get("format", "doc")
            text_content = doc_item["raw_text"]

            # Формируем структурированный документ RAG
            docs_to_ingest.append(Document(
                doc_id=f"client_doc_{company_name}_{idx+1}",
                text=(
                    f"Документ клиента '{file_name}' (Формат: {fmt.upper()}, Компания: {company_name}, Ниша: {niche}):\n"
                    f"{text_content}"
                ),
                metadata={
                    "category": "client_files",
                    "company_name": company_name,
                    "file_name": file_name,
                    "format": fmt
                }
            ))

        if docs_to_ingest:
            indexed_count = await rag_pipeline.ingest_documents_async(docs_to_ingest)
            logger.info(f"[DocumentCollector] 📚 Успешно заиндексировано {indexed_count} чанков из {len(docs_to_ingest)} клиентских файлов в RAG для {company_name}.")
            return indexed_count
        return 0

    def synthesize_dossier_from_docs(self, extracted_docs: List[Dict[str, Any]]) -> str:
        """
        Формирует единое текстовое досье из всех загруженных файлов.
        """
        valid_texts = [
            f"=== Файл: {d['file_name']} ({d['format'].upper()}) ===\n{d['raw_text']}"
            for d in extracted_docs
            if d.get("status") == "success" and d.get("raw_text")
        ]
        return "\n\n".join(valid_texts)
